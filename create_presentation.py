# -*- coding: utf-8 -*-
"""Vaveyla bitirme projesi sunumu oluşturucu."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = r"c:\Users\erdem\Desktop\Vaveyla_Sunum.pptx"

# Renk paleti — pastane temasına uygun sıcak tonlar
PRIMARY = RGBColor(0x8B, 0x1A, 0x4A)      # bordo
SECONDARY = RGBColor(0xD4, 0x6B, 0x8A)     # pembe
ACCENT = RGBColor(0xF5, 0xE6, 0xD3)        # krem
DARK = RGBColor(0x2C, 0x2C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF7, 0xF7, 0xF7)
MUTED = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header_bar(slide, title_text):
    add_rect(slide, 0, 0, W, Inches(1.15), PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), W - Inches(1.2), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def add_bullets(slide, items, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), font_size=22):
    box = slide.shapes.add_textbox(left, top, width, H - top - Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        p.line_spacing = 1.25


def add_two_column(slide, left_title, left_items, right_title, right_items):
    # Sol sütun başlık
    lt = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(5.8), Inches(0.5))
    lp = lt.text_frame.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(20)
    lp.font.bold = True
    lp.font.color.rgb = PRIMARY
    lp.font.name = "Calibri"

    lb = slide.shapes.add_textbox(Inches(0.6), Inches(1.85), Inches(5.8), Inches(5.0))
    ltf = lb.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(6)

    # Sağ sütun başlık
    rt = slide.shapes.add_textbox(Inches(6.9), Inches(1.35), Inches(5.8), Inches(0.5))
    rp = rt.text_frame.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(20)
    rp.font.bold = True
    rp.font.color.rgb = PRIMARY
    rp.font.name = "Calibri"

    rb = slide.shapes.add_textbox(Inches(6.9), Inches(1.85), Inches(5.8), Inches(5.0))
    rtf = rb.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(6)

    add_rect(slide, Inches(6.55), Inches(1.35), Inches(0.04), Inches(5.5), SECONDARY)


# ── SLIDE 1: Kapak ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_rect(slide, 0, Inches(2.8), W, Inches(0.06), ACCENT)

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), W - Inches(1.6), Inches(1.2))
tp = title_box.text_frame.paragraphs[0]
tp.text = "VAVEYLA"
tp.font.size = Pt(54)
tp.font.bold = True
tp.font.color.rgb = WHITE
tp.font.name = "Calibri"
tp.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), W - Inches(1.6), Inches(1.0))
sp = sub_box.text_frame.paragraphs[0]
sp.text = "Pastanelere Özel Çok Panelli Sipariş ve Yönetim Sistemi"
sp.font.size = Pt(24)
sp.font.color.rgb = ACCENT
sp.font.name = "Calibri"
sp.alignment = PP_ALIGN.CENTER

info_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), W - Inches(1.6), Inches(2.5))
itf = info_box.text_frame
itf.word_wrap = True
lines = [
    "Trakya Üniversitesi · Bilgisayar Mühendisliği · BLM421 Proje II",
    "",
    "İrem Su ERDEMİR  ·  Beyza YILMAZ",
    "Proje Danışmanı: Dr. Öğr. Üyesi Andaç MESUT",
    "Edirne · 2026",
]
for i, line in enumerate(lines):
    p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)


# ── SLIDE 2: Projenin Amacı ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Projenin Amacı")
add_bullets(slide, [
    "Pastanelere özel, çok panelli bir sipariş ve yönetim platformu geliştirmek",
    "Müşteri, pastane, kurye ve admin rollerinin aynı sistem üzerinde etkileşim kurmasını sağlamak",
    "Sipariş süreçlerini dijitalleştirerek hızlı, düzenli ve erişilebilir hale getirmek",
    "Kullanıcı dostu arayüz ile pastane ürünlerine (pasta, tatlı, unlu mamul) kolay erişim sunmak",
    "Modüler mimari ile ileride genişletilebilir bir altyapı oluşturmak",
])


# ── SLIDE 3: Çözülen Problemler ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Çözülen Problemler")
add_bullets(slide, [
    "Pastaneler için dijital sipariş sisteminin eksikliği",
    "Müşterilerin özel ürünlere (pasta, tatlı vb.) kolay erişememesi",
    "Sipariş ve teslimat süreçlerinin takibinde yaşanan zorluklar",
    "Kurye, müşteri ve işletme arasındaki iletişim kopuklukları",
    "Merkezi yönetim ve denetim eksikliği",
], font_size=21)


# ── SLIDE 4: Sistem Mimarisi ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Sistem Mimarisi — 4 Panel")

panels = [
    ("Müşteri Paneli", "Ürün görüntüleme\nSipariş verme\nSipariş takibi\nKurye ile iletişim"),
    ("Pastane Paneli", "Ürün yönetimi\nSipariş yönetimi\nMenü & stok\nKampanya yönetimi"),
    ("Kurye Paneli", "Atanan siparişler\nTeslimat süreci\nDurum güncelleme\nMüşteri iletişimi"),
    ("Admin Paneli", "Kullanıcı yönetimi\nSistem denetimi\nSipariş izleme\nRaporlama"),
]

panel_w = Inches(2.85)
gap = Inches(0.25)
start_x = Inches(0.55)
y_top = Inches(1.6)

for idx, (title, desc) in enumerate(panels):
    x = start_x + idx * (panel_w + gap)
    card = add_rect(slide, x, y_top, panel_w, Inches(4.8), LIGHT_GRAY)
    # Üst şerit
    add_rect(slide, x, y_top, panel_w, Inches(0.65), SECONDARY)

    tb = slide.shapes.add_textbox(x + Inches(0.15), y_top + Inches(0.1), panel_w - Inches(0.3), Inches(0.5))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(17)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.font.name = "Calibri"
    tp.alignment = PP_ALIGN.CENTER

    db = slide.shapes.add_textbox(x + Inches(0.2), y_top + Inches(0.85), panel_w - Inches(0.4), Inches(3.8))
    dtf = db.text_frame
    dtf.word_wrap = True
    for j, line in enumerate(desc.split("\n")):
        p = dtf.paragraphs[0] if j == 0 else dtf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(8)

# Alt mimari notu
note = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), W - Inches(1.2), Inches(0.7))
np = note.text_frame.paragraphs[0]
np.text = "Flutter (Cross-Platform Frontend)  ←→  .NET 8 REST API  ←→  MS SQL Server"
np.font.size = Pt(16)
np.font.bold = True
np.font.color.rgb = PRIMARY
np.font.name = "Calibri"
np.alignment = PP_ALIGN.CENTER


# ── SLIDE 5: Kullanıcı Rolleri ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Hedef Kitle ve Kullanıcı Rolleri")
add_two_column(
    slide,
    "Son Kullanıcılar",
    [
        "Müşteriler — sipariş veren son kullanıcılar",
        "Kuryeler — teslimat sürecini yürüten kullanıcılar",
    ],
    "İşletme & Yönetim",
    [
        "Pastaneler — ürün ve sipariş yönetimi yapan işletmeler",
        "Admin — sistem genelinde denetim ve yönetim",
    ],
)

# Platform notu
plat = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), W - Inches(1.2), Inches(0.8))
pp = plat.text_frame.paragraphs[0]
pp.text = "Platform Desteği: Android · iOS · Web (Windows, macOS, Linux)"
pp.font.size = Pt(16)
pp.font.color.rgb = MUTED
pp.font.name = "Calibri"
pp.alignment = PP_ALIGN.CENTER


# ── SLIDE 6: Frontend Teknolojileri ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Kullanılan Teknolojiler — Frontend")
add_two_column(
    slide,
    "Temel",
    [
        "Flutter & Dart — cross-platform UI",
        "BLoC — durum yönetimi",
        "Material Design — arayüz bileşenleri",
        "easy_localization — çoklu dil desteği",
    ],
    "Kütüphaneler",
    [
        "flutter_map — harita ve konum",
        "signalr_netcore — gerçek zamanlı iletişim",
        "http — REST API istekleri",
        "shared_preferences — yerel veri depolama",
        "flutter_local_notifications — bildirimler",
        "image_picker — görsel yükleme",
    ],
)


# ── SLIDE 7: Backend Teknolojileri ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Kullanılan Teknolojiler — Backend")
add_two_column(
    slide,
    "Sunucu & Veritabanı",
    [
        "C# / .NET 8 Web API",
        "Entity Framework Core",
        "MS SQL Server",
        "Dapper — performanslı sorgular",
        "RESTful API mimarisi",
    ],
    "Güvenlik & Araçlar",
    [
        "JWT Bearer Authentication",
        "BCrypt — şifre hashleme",
        "MailKit — e-posta (şifre sıfırlama)",
        "Swagger / OpenAPI — API dokümantasyonu",
        "Visual Studio — geliştirme ortamı",
        "Git & GitHub — sürüm kontrolü",
    ],
)


# ── SLIDE 8: Temel Modüller ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Temel Modüller ve Özellikler")
add_bullets(slide, [
    "Kullanıcı yönetimi — kayıt, giriş, şifre sıfırlama, rol bazlı erişim",
    "Ürün & kategori yönetimi — menü, stok, kampanya ve kupon sistemi",
    "Sipariş süreçleri — sepet, sipariş oluşturma, durum takibi",
    "Gerçek zamanlı sohbet — müşteri ↔ kurye / pastane iletişimi (SignalR)",
    "Harita entegrasyonu — yakındaki pastaneleri görüntüleme",
    "Bildirim sistemi — sipariş ve kampanya bildirimleri",
    "Admin paneli — kullanıcı denetimi, sipariş izleme, raporlama",
], font_size=19)


# ── SLIDE 9: Farklılaştırıcı Özellikler ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Vaveyla'nın Farklılaştırıcı Özellikleri")
add_bullets(slide, [
    "Pastanelere özel niş odak — Yemeksepeti/GetirYemek gibi genel platformlardan ayrışma",
    "4 ayrı panel — müşteri, pastane, kurye ve admin için rol bazlı detaylı yönetim",
    "Kurye ile doğrudan iletişim — teslimat sürecinde şeffaf etkileşim",
    "Merkezi admin denetimi — sistem güvenliği ve sürdürülebilirlik",
    "Modüler & genişletilebilir mimari — online ödeme, AI önerileri vb. için hazır altyapı",
    "Cross-platform — tek kod tabanı ile mobil ve web desteği",
], font_size=19)


# ── SLIDE 10: Test ve Geri Bildirim ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Test ve Kullanıcı Geri Bildirimi")
add_two_column(
    slide,
    "Test Süreci",
    [
        "Fonksiyonel test senaryoları uygulandı",
        "Android emülatör ve fiziksel cihazda test",
        "Backend API — Postman & Swagger ile doğrulama",
        "Windows 11 · i5 · 8 GB RAM ortamında test",
    ],
    "Anket Sonuçları (20 katılımcı)",
    [
        "Genel memnuniyet: 4,15 / 5",
        "Estetik değerlendirme: 4,25 / 5",
        "NPS (Net Promoter Score): +40",
        "Kullanıcı arayüzü: 4,35 / 5",
        "%90 restoran/ürün listeleme kullanımı",
    ],
)


# ── SLIDE 11: Sonuç ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Sonuç")
add_bullets(slide, [
    "Vaveyla, pastane ürünlerinin dijital ortamda satışını destekleyen bütünleşik bir sistemdir",
    "4 farklı kullanıcı rolü için ayrı paneller ile kapsamlı yönetim sunulmaktadır",
    "Flutter + .NET 8 + SQL Server mimarisi ile güvenilir ve ölçeklenebilir altyapı sağlanmıştır",
    "Kullanıcı geri bildirimleri olumlu — sistem ihtiyaçları büyük ölçüde karşılamaktadır",
    "Modüler yapı sayesinde gelecekte yeni özellikler kolayca eklenebilir",
], font_size=20)


# ── SLIDE 12: Gelecek Planları ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Gelecek Geliştirmeler")
add_bullets(slide, [
    "Kişiye özel pasta tasarımı ve yapay zeka destekli ürün önerileri",
    "Sadakat / puan sistemi ve gelişmiş kampanya yönetimi",
    "Çevrimiçi ödeme entegrasyonu (kredi kartı, dijital cüzdan, QR)",
    "Canlı sipariş takibi ve gelişmiş bildirim sistemi",
    "Veri analizi & raporlama modülü",
    "Çok dilli destek genişletmesi",
], font_size=20)


# ── SLIDE 13: Teşekkürler ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_rect(slide, 0, Inches(3.4), W, Inches(0.06), ACCENT)

tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), W - Inches(1.6), Inches(1.2))
tp = tb.text_frame.paragraphs[0]
tp.text = "Teşekkürler"
tp.font.size = Pt(48)
tp.font.bold = True
tp.font.color.rgb = WHITE
tp.font.name = "Calibri"
tp.alignment = PP_ALIGN.CENTER

sb = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), W - Inches(1.6), Inches(1.5))
stf = sb.text_frame
for i, line in enumerate([
    "Sorularınız için hazırız.",
    "",
    "Canlı demo ile devam edeceğiz.",
]):
    p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)

prs.save(OUTPUT)
print(f"Sunum oluşturuldu: {OUTPUT}")
print(f"Toplam slayt: {len(prs.slides)}")
